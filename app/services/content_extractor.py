"""
Comprehensive Content Extraction Service

Extracts detailed study information from multiple file types:
- Topics & subtopics
- Learning objectives
- Key terms & definitions
- Complexity analysis
- Study time estimates
"""
import base64
import json
from typing import List, Tuple
from io import BytesIO

from app.models import StudyMaterial, Topic, KeyTerm
from app.clients.ai_client import get_ai_client


def extract_comprehensive_content(
    filename: str,
    file_content: str,
    file_type: str,
    user_id: str
) -> Tuple[StudyMaterial, List[str]]:
    """
    Extract comprehensive study content from a file

    Args:
        filename: Original filename
        file_content: Base64 encoded content
        file_type: Type of file (pdf, image, pptx, etc.)
        user_id: User ID

    Returns:
        Tuple of (StudyMaterial, warnings)
    """
    warnings = []

    # Decode file
    try:
        decoded = base64.b64decode(file_content)
    except Exception as e:
        raise ValueError(f"Invalid file content: {str(e)}")

    # Extract text/images based on type
    if file_type.lower() in ['pdf']:
        text, images = _extract_from_pdf(decoded)
    elif file_type.lower() in ['jpg', 'jpeg', 'png', 'gif']:
        text = None
        images = [file_content]  # Already base64
    elif file_type.lower() in ['pptx', 'ppt']:
        text, images = _extract_from_pptx(decoded)
        warnings.append("PowerPoint extraction may be limited - some formatting lost")
    else:
        # Plain text
        text = decoded.decode('utf-8')
        images = None

    # Use AI to extract comprehensive content
    ai_client = get_ai_client()

    try:
        extraction_result = ai_client.extract_comprehensive_study_content(
            text=text,
            images=images,
            filename=filename
        )
    except Exception as e:
        raise Exception(f"AI extraction failed: {str(e)}")

    # Build StudyMaterial object
    topics = []
    total_hours = 0.0

    for topic_data in extraction_result.get("topics", []):
        # Extract key terms
        key_terms = [
            KeyTerm(
                term=kt["term"],
                definition=kt["definition"],
                importance=kt.get("importance", "medium")
            )
            for kt in topic_data.get("key_terms", [])
        ]

        # Create Topic
        topic = Topic(
            name=topic_data["name"],
            subtopics=topic_data.get("subtopics", []),
            complexity=topic_data.get("complexity", "intermediate"),
            estimated_hours=topic_data.get("estimated_hours", 3.0),
            learning_objectives=topic_data.get("learning_objectives", []),
            key_terms=key_terms
        )

        topics.append(topic)
        total_hours += topic.estimated_hours

    # Create StudyMaterial
    study_material = StudyMaterial(
        user_id=user_id,
        source_file=filename,
        file_type=file_type,
        topics=topics,
        total_estimated_hours=total_hours,
        complexity_level=extraction_result.get("overall_complexity", "mixed"),
        raw_content=text  # Store full extracted text for tutor
    )

    return study_material, warnings


def _extract_from_pdf(pdf_bytes: bytes) -> Tuple[str, List[str]]:
    """
    Extract text from PDF using OCR and text extraction

    Cost-optimized approach:
    1. Try PyPDF2 text extraction first (free, fast)
    2. If that fails or returns little text, use OCR (free, slower)
    3. Send only TEXT to AI (cheap model like gpt-4o-mini)
    """
    from PyPDF2 import PdfReader

    # First, try native text extraction
    try:
        pdf_reader = PdfReader(BytesIO(pdf_bytes))
        text = ""
        total_pages = len(pdf_reader.pages)

        print(f"PDF has {total_pages} pages, extracting text...")

        for page in pdf_reader.pages:  # Process ALL pages
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n\n"

        text = text.strip()

        # If we got decent text (more than 100 chars), use it
        if len(text) > 100:
            print(f"PDF text extraction successful: {len(text)} chars")
            return text, None  # No images needed - text-only extraction

        print(f"PDF text extraction yielded little text ({len(text)} chars), trying OCR...")
    except Exception as e:
        print(f"PDF text extraction failed: {e}, trying OCR...")

    # Fallback: OCR for scanned PDFs
    try:
        from pdf2image import convert_from_bytes
        import pytesseract

        # Convert PDF to images
        images = convert_from_bytes(pdf_bytes, dpi=200)  # Higher DPI for better OCR

        ocr_text = ""
        for i, img in enumerate(images[:20]):  # Limit to 20 pages
            print(f"OCR processing page {i+1}/{min(len(images), 20)}...")
            page_text = pytesseract.image_to_string(img)
            ocr_text += f"\n\n--- Page {i+1} ---\n{page_text}"

        print(f"OCR extraction successful: {len(ocr_text)} chars")
        return ocr_text.strip(), None  # Return text only, no images

    except ImportError as e:
        print(f"OCR libraries not available: {e}")
        print("Install with: pip install pytesseract pdf2image")
        print("Also install Tesseract OCR: https://github.com/tesseract-ocr/tesseract")

        # Return whatever text we got
        return text if text else "Could not extract text from PDF", None
    except Exception as e:
        print(f"OCR failed: {e}")
        return text if text else "Could not extract text from PDF", None


def _extract_from_pptx(pptx_bytes: bytes) -> Tuple[str, List[str]]:
    """Extract text from PowerPoint (basic)"""
    try:
        from pptx import Presentation

        prs = Presentation(BytesIO(pptx_bytes))
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text.strip(), None

    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    except Exception as e:
        raise ValueError(f"Failed to parse PowerPoint: {str(e)}")
